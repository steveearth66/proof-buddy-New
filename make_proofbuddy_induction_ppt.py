from pptx import Presentation
from pptx.util import Inches, Pt

# Monospace font available on Windows
CODE_FONT = "Consolas"

SLIDES = [
    {
        "title": "Proof Buddy Backend Induction: Informal Correctness",
        "bullets": [
            "Claim: backend rejects malformed / ill-typed / unsound steps",
            "Approach: map each error class → exact code check → unit tests",
        ],
        "code": "",
        "notes": """Say: 'I’ll summarize an informal proof of correctness for the backend induction process.'
Say: 'Each step is gated by parse + typing + rule applicability; invalid steps fail fast with clear errors.'"""
    },
    {
        "title": "Backend Pipeline (one proof step)",
        "bullets": [
            "Input: expression + rule string",
            "Pipeline: parse → label → decorate → type/arity checks → apply rule",
            "Reject with specific error messages",
        ],
        "code": """# ERProofLine.applyRule (excerpt)
parts = rule.split()
ruleCategory = parts[0] if parts else ""
rule = parts[1] if len(parts) > 1 else ""

if ruleCategory not in ("eval", "apply", "rewrite"):
    self.errLog.append("Rule must start with 'eval', 'apply', or 'rewrite'")
    return

for label in self.find_undefined_labels(targetNode):
    self.errLog.append(f"No definition found for label '{label}'")""",
        "notes": """Open ERProofEngine.applyRule.
Emphasize: global gatekeeping happens before any rewrite/eval."""
    },
    {
        "title": "Why Racket-lite + list processing",
        "bullets": [
            "Core for recursion: cons/first/rest/null? + if",
            "Enough arithmetic/logic to express typical induction goals",
            "Safety: small surface area → easier to typecheck and test",
        ],
        "code": """# Labeler.py builtins (excerpt)
BUILT_IN_FUNCTIONS = [
  'if','cons','first','rest','null?',
  '+','-','*','quotient','remainder','zero?',
  'expt','=','<=','>=','<','>','and','or','not','xor','implies',
  'list?','integer?'
]""",
        "notes": """Say: 'This is the minimal core to express structural recursion + common arithmetic invariants.'
Point out: limiting the language makes sound checking tractable."""
    },
    {
        "title": "Name Resolution: builtin vs UDF vs generic vs undefined",
        "bullets": [
            "Builtin → fixed type + fixed arity",
            "If not builtin but in defDict → treat as UDF",
            "If label appears but has no definition → backend error",
        ],
        "code": """# Labeler.labelTree (excerpt)
if inputTree.data in BUILT_IN_FUNCTIONS:
    erObj = pdict[inputTree.data]
    inputTree.type = RacType((tuple([(None, t) for t in erObj.ins]), (None, erObj.outtype)))
    inputTree.numArgs = erObj.numArgs
elif inputTree.data in defDict:
    inputTree.type = defDict[inputTree.data].racType
    if inputTree.type.isType("FUNCTION"):
        inputTree.numArgs = len(inputTree.type.getDomain())""",
        "notes": """Say: 'This is where we decide what a symbol *means* before we typecheck.'
Mention: undefined labels are caught during applyRule (previous slide)."""
    },
    {
        "title": "Arity Checking (Wrong # of arguments)",
        "bullets": [
            "Every call node checks expectedCount vs providedCount",
            "Arity fails before deeper rule checks",
            "Unit tests cover many arity failures (e.g., +, -, *, expt)",
        ],
        "code": """# Decorator.argQty (excerpt)
expectedCount = func.numArgs
providedCount = len(treeNode.children) - 1
if (expectedCount != None) and (expectedCount != providedCount):
    return [False,
      f"{func.name} only takes {expectedCount} arguments, but {providedCount} were provided"
    ]""",
        "notes": """Open Decorator.argQty.
Then show: test_math_operations has '+ only takes 2 arguments...' cases."""
    },
    {
        "title": "Type Checking (Domain match) + if typing",
        "bullets": [
            "Call typecheck compares provided arg types vs expected domain types",
            "if: guard must be BOOL; branches must have matching types",
            "These checks block nonsense steps early",
        ],
        "code": """# Decorator.typeCheck (excerpt)
providedIns = [c.type for c in inputTree.children[1:]]
expectedIns = [(RacType(x) if isinstance(x,tuple) else x) for x in func.type.value[0]]
if not all(x==y for x, y in zip(providedIns, expectedIns)):
    return [False, "Cannot match argument out typeList ..."]

# Decorator.remTemps (excerpt)
if func.data == "if":
    if not inputTree.children[1].type.isType("BOOL"):
        errLog.append("The first argument of an if function must be Boolean ...")""",
        "notes": """Say: 'We enforce well-typedness of the expression language, independent of the proof rules.'
Point out: if typing is special because output depends on branches."""
    },
    {
        "title": "Applying a UDF: Assignment Validation (names/types/values)",
        "bullets": [
            "Checks x=... syntax, count, name/order",
            "Typechecks each assignment value",
            "Verifies assigned values match the selected node’s actual args",
        ],
        "code": """# ERProofEngine.parse_and_typecheck_args (excerpt)
for param in rawParams:
    if '=' not in param:
        self.errLog.append("... Did you forget an equals sign?")
        return [], True
    elif param.count('=') > 1:
        self.errLog.append("... Did you forget a comma?")
        return [], True

if got < need: self.errLog.append("Not enough arguments ..."); return [], True
if got > need: self.errLog.append("Too many arguments ..."); return [], True

# name/order check
name, _ = param.split('=', 1)
if name.strip() != expected:
    self.errLog.append("Argument ... is in position ... but expected ...")

# type mismatch
if typed.type != expected_type.getType():
    self.errLog.append("Type mismatch in argument ...")

# value mismatch vs target node
if user_val != target_val:
    self.errLog.append("Value mismatch in argument ...")""",
        "notes": """Emphasize: this corresponds directly to advisor checklist:
name, arity, param types, output consistency via typing, and exact matching behavior."""
    },
    {
        "title": "Unit Tests: UDF Positives + Negatives",
        "bullets": [
            "Wrong prefix rejected",
            "Wrong arity / comma mistakes rejected",
            "Wrong param name/order/type/value rejected",
            "Correct case passes",
        ],
        "code": """# proofs/test_axioms_and_udfs.py (excerpt)
do_single_test_case("fc", "(fc 3 4)",
  ["Rule must start with 'eval', 'apply', or 'rewrite'"], udfProof)

do_single_test_case("apply fc x=3 y=4", "(fc 3 4)",
  ["Too many assignments for a given argument 'x=3 y=4'. Did you forget a comma?"], udfProof)

do_single_test_case("apply fc z=3, y=4", "(fc 3 4)",
  ["Argument 'z' is in position 1 but expected 'x' for fc"], udfProof)

do_single_test_case("apply fc x=#t, y=4", "(fc 3 4)",
  ["Type mismatch in argument 'x=#t': expected INT, got BOOL"], udfProof)

do_single_test_case("apply fc x=3, y=4", "(fc 3 4)", "(* 3 4)", udfProof)""",
        "notes": """Say: 'These tests are essentially a spec for our error handling.'
If asked: 'Why so many negative cases?' → Because most student mistakes are format/name/type mistakes."""
    },
    {
        "title": "IH Rule: Exact Match Only",
        "bullets": [
            "IH takes no parameters",
            "Applies only if selected node exactly matches IH LHS or IH RHS",
            "Prevents unsound “IH anywhere” rewrites",
        ],
        "code": """# ERRuleset.IH.isApplicable (excerpt)
if rawParams:
    return False, "IH rule takes no parameters"

nodeStr = str(ruleNode)
lhsStr = str(self.indHypLHS)
rhsStr = str(self.indHypRHS)

if nodeStr == lhsStr or nodeStr == rhsStr:
    return True, "IH.isApplicable() PASS"
return False, f"Node '{nodeStr}' does not match ...\"""",
        "notes": """Call out the key correctness point: IH only fires when the user selects exactly the hypothesis expression."""
    },
    {
        "title": "Builtin arity/type tests (quick example)",
        "bullets": [
            "Builtins have strong unit test coverage",
            "Example: + arity/type errors and a passing case",
        ],
        "code": """# proofs/test_math_operations.py (excerpt)
("(+ 1)", ['+ only takes 2 arguments, but 1 was provided']),
("(+ 1 1 1)", ['+ only takes 2 arguments, but 3 were provided']),
("(+ 1 #t)", ["Cannot match argument out typeList ['INT', 'BOOL'] with expected typeList ['INT', 'INT']"]),
("(+ 1 2)", 3)""",
        "notes": """Say: 'This shows the same “early rejection” behavior for the expression language itself.'"""
    },
    {
        "title": "Conclusion: Informal Correctness Argument",
        "bullets": [
            "Accepted steps are: well-formed, well-typed, rule-justified",
            "Backend blocks key error classes with explicit messages",
            "Limitations: IH is string-match; expressiveness bounded by ruleset",
        ],
        "code": "",
        "notes": """Close with: 'Correctness here means we never accept a step that violates these syntactic/type/rule constraints.'
If asked about limitations: mention IH uses string equality, not semantic equivalence."""
    },
]

def _set_paragraph_font(p, size_pt, font_name=None, bold=None):
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    r.font.size = Pt(size_pt)
    if font_name:
        r.font.name = font_name
    if bold is not None:
        r.font.bold = bold

def add_bullets(slide, title, bullets):
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        _set_paragraph_font(p, 22)

def add_code_box(slide, code_text):
    if not code_text.strip():
        return

    # Place a code box under the bullet content area.
    left = Inches(0.7)
    top = Inches(4.0)
    width = Inches(12.0)
    height = Inches(3.0)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()

    header = tf.paragraphs[0]
    header.text = "Code snippet"
    _set_paragraph_font(header, 14, font_name=CODE_FONT, bold=True)

    for line in code_text.strip("\n").splitlines():
        p = tf.add_paragraph()
        p.text = line.rstrip("\n")
        p.level = 0
        _set_paragraph_font(p, 12, font_name=CODE_FONT)

def add_notes(slide, notes_text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = notes_text

def main():
    prs = Presentation()
    for s in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        add_bullets(slide, s["title"], s["bullets"])
        add_code_box(slide, s.get("code", ""))
        add_notes(slide, s["notes"])

    prs.save("ProofBuddy_Backend_Induction_Correctness.pptx")

if __name__ == "__main__":
    main()